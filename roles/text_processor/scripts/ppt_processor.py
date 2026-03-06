#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PowerPoint 演示文稿处理器
提供 PPT 创建、编辑、图表、动画等功能
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
from pathlib import Path


class PPTProcessor:
    """PowerPoint 演示文稿处理类"""
    
    def __init__(self):
        self.prs = None
        self.slide = None
    
    # ==================== 基础操作 ====================
    
    def create(self):
        """创建新演示文稿"""
        self.prs = Presentation()
        return self
    
    def load(self, filepath):
        """加载现有演示文稿"""
        self.prs = Presentation(filepath)
        return self
    
    def save(self, filepath):
        """保存演示文稿"""
        self.prs.save(filepath)
        print(f"✅ 已保存: {filepath}")
        return filepath
    
    # ==================== 幻灯片操作 ====================
    
    def add_title_slide(self, title, subtitle=""):
        """添加标题幻灯片"""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if subtitle:
            slide.placeholders[1].text = subtitle
        self.slide = slide
        return slide
    
    def add_content_slide(self, title, content=""):
        """添加内容幻灯片"""
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if content:
            slide.placeholders[1].text = content
        self.slide = slide
        return slide
    
    def add_blank_slide(self):
        """添加空白幻灯片"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide = slide
        return slide
    
    def use_layout(self, layout_index):
        """使用指定布局"""
        slide_layout = self.prs.slide_layouts[layout_index]
        slide = self.prs.slides.add_slide(slide_layout)
        self.slide = slide
        return slide
    
    # ==================== 文本操作 ====================
    
    def add_text(self, text, left, top, width, height, font_size=12, bold=False):
        """添加文本框"""
        txBox = self.slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text
        
        # 设置字体
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = bold
        
        return txBox
    
    def add_bullet_text(self, bullets, level=0):
        """添加项目符号文本"""
        tf = self.slide.placeholders[1].text_frame
        tf.text = bullets[0]
        
        for bullet in bullets[1:]:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = level
        
        return tf
    
    def add_numbered_list(self, items):
        """添加编号列表"""
        tf = self.slide.placeholders[1].text_frame
        tf.text = items[0]
        
        for i, item in enumerate(items[1:], 1):
            p = tf.add_paragraph()
            p.text = f"{i}. {item}"
        
        return tf
    
    # ==================== 图表操作 ====================
    
    def add_chart(self, title, chart_type, data):
        """添加图表"""
        chart_data = CategoryChartData()
        chart_data.categories = data['categories']
        for series_name, series_data in data['series'].items():
            chart_data.add_series(series_name, series_data)
        
        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)
        chart = self.slide.shapes.add_chart(
            chart_type, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_title = True
        chart.chart_title.text_frame.text = title
        
        return chart
    
    def add_column_chart(self, title, categories, series_data):
        """添加柱状图"""
        data = {
            'categories': categories,
            'series': series_data
        }
        return self.add_chart(title, XL_CHART_TYPE.COLUMN_CLUSTERED, data)
    
    def add_line_chart(self, title, categories, series_data):
        """添加折线图"""
        data = {
            'categories': categories,
            'series': series_data
        }
        return self.add_chart(title, XL_CHART_TYPE.LINE, data)
    
    def add_pie_chart(self, title, categories, series_data):
        """添加饼图"""
        data = {
            'categories': categories,
            'series': series_data
        }
        return self.add_chart(title, XL_CHART_TYPE.PIE, data)
    
    def add_bar_chart(self, title, categories, series_data):
        """添加条形图"""
        data = {
            'categories': categories,
            'series': series_data
        }
        return self.add_chart(title, XL_CHART_TYPE.BAR_CLUSTERED, data)
    
    # ==================== 表格操作 ====================
    
    def add_table(self, data, left=None, top=None, width=None, height=None):
        """添加表格"""
        rows = len(data)
        cols = len(data[0]) if data else 0
        
        if left is None:
            left = Inches(1)
        if top is None:
            top = Inches(2)
        if width is None:
            width = Inches(8)
        
        table = self.slide.shapes.add_table(rows, cols, left, top, width, height).table
        table.style = 'Table Grid'
        
        for row_idx, row_data in enumerate(data):
            for col_idx, value in enumerate(row_data):
                table.cell(row_idx, col_idx).text = str(value)
        
        return table
    
    def add_simple_table(self, data, columns):
        """添加简单表格（带表头）"""
        header = list(columns.values())
        rows_data = [header]
        for row in data:
            rows_data.append([str(row.get(col, '')) for col in columns.keys()])
        return self.add_table(rows_data)
    
    # ==================== 图片和形状 ====================
    
    def add_image(self, image_path, width=None, height=None, left=None, top=None):
        """添加图片"""
        if left is None:
            left = Inches(1)
        if top is None:
            top = Inches(2)
        
        if width:
            self.slide.shapes.add_picture(image_path, left, top, width=Inches(width))
        elif height:
            self.slide.shapes.add_picture(image_path, left, top, height=Inches(height))
        else:
            self.slide.shapes.add_picture(image_path, left, top)
        
        return self
    
    def add_shape(self, shape_type, left, top, width, height, text=""):
        """添加形状"""
        shape = self.slide.shapes.add_shape(
            shape_type, left, top, width, height
        )
        if text:
            shape.text = text
        return shape
    
    def add_rectangle(self, text="", left=Inches(1), top=Inches(2), width=Inches(2), height=Inches(1)):
        """添加矩形"""
        return self.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height, text)
    
    def add_oval(self, text="", left=Inches(1), top=Inches(2), width=Inches(2), height=Inches(1)):
        """添加椭圆"""
        return self.add_shape(MSO_SHAPE.OVAL, left, top, width, height, text)
    
    def add_arrow(self, left, top, width, height, direction='right'):
        """添加箭头"""
        if direction == 'right':
            shape_type = MSO_SHAPE.RIGHT_ARROW
        elif direction == 'left':
            shape_type = MSO_SHAPE.LEFT_ARROW
        elif direction == 'up':
            shape_type = MSO_SHAPE.UP_ARROW
        else:
            shape_type = MSO_SHAPE.DOWN_ARROW
        
        return self.add_shape(shape_type, left, top, width, height)
    
    # ==================== 模板生成 ====================
    
    @staticmethod
    def create_weekly_report(filepath, week_info):
        """创建周报 PPT"""
        prs = Presentation()
        
        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = f"{week_info['project']} - 周报"
        slide.placeholders[1].text = f"{week_info['date_range']}"
        
        # 目录
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "目录"
        tf = slide.placeholders[1].text_frame
        sections = ['本周完成', '进行中', '问题与风险', '下周计划']
        for i, section in enumerate(sections):
            p = tf.add_paragraph()
            p.text = f"{i+1}. {section}"
        
        # 本周完成
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "本周完成"
        tf = slide.placeholders[1].text_frame
        for task in week_info.get('completed', []):
            p = tf.add_paragraph()
            p.text = f"✅ {task}"
        
        # 进行中
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "进行中"
        tf = slide.placeholders[1].text_frame
        for task in week_info.get('in_progress', []):
            p = tf.add_paragraph()
            p.text = f"🔄 {task}"
        
        # 问题与风险
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "问题与风险"
        tf = slide.placeholders[1].text_frame
        for risk in week_info.get('risks', []):
            p = tf.add_paragraph()
            p.text = f"⚠️ {risk}"
        
        # 下周计划
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "下周计划"
        tf = slide.placeholders[1].text_frame
        for i, plan in enumerate(week_info.get('next_week', [])):
            p = tf.add_paragraph()
            p.text = f"📋 {plan}"
        
        # 进度图表
        if 'progress' in week_info:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "整体进度"
            
            chart_data = CategoryChartData()
            chart_data.categories = ['已完成', '进行中', '待开始']
            chart_data.add_series('任务数', week_info['progress'])
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(1), Inches(2), Inches(8), Inches(4.5), chart_data
            ).chart
        
        prs.save(filepath)
        print(f"✅ 周报 PPT 已创建: {filepath}")
        return filepath
    
    @staticmethod
    def create_project_report(filepath, project_info):
        """创建项目报告 PPT"""
        prs = Presentation()
        
        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = project_info['title']
        slide.placeholders[1].text = f"项目经理: {project_info['manager']}\n日期: {project_info['date']}"
        
        # 项目概览
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "项目概览"
        tf = slide.placeholders[1].text_frame
        tf.text = f"项目名称: {project_info['name']}"
        p = tf.add_paragraph()
        p.text = f"当前阶段: {project_info['phase']}"
        p = tf.add_paragraph()
        p.text = f"整体进度: {project_info['progress']}%"
        p = tf.add_paragraph()
        p.text = f"团队规模: {project_info['team_size']}人"
        
        # 进度统计
        if 'task_stats' in project_info:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = "任务统计"
            
            chart_data = CategoryChartData()
            chart_data.categories = list(project_info['task_stats'].keys())
            chart_data.add_series('任务数', list(project_info['task_stats'].values()))
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1), Inches(2), Inches(8), Inches(4.5), chart_data
            ).chart
        
        # 风险与问题
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "问题与风险"
        tf = slide.placeholders[1].text_frame
        for risk in project_info.get('risks', []):
            p = tf.add_paragraph()
            p.text = f"⚠️ {risk}"
        
        # 下一步计划
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "下一步计划"
        tf = slide.placeholders[1].text_frame
        for i, plan in enumerate(project_info.get('plans', [])):
            p = tf.add_paragraph()
            p.text = f"{i+1}. {plan}"
        
        prs.save(filepath)
        print(f"✅ 项目报告 PPT 已创建: {filepath}")
        return filepath
    
    @staticmethod
    def create_product_demo(filepath, product_info):
        """创建产品演示 PPT"""
        prs = Presentation()
        
        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = product_info['name']
        slide.placeholders[1].text = product_info.get('tagline', '')
        
        # 目录
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "目录"
        tf = slide.placeholders[1].text_frame
        for i, section in enumerate(product_info['sections']):
            p = tf.add_paragraph()
            p.text = f"{i+1}. {section}"
        
        # 特性介绍
        for i, feature in enumerate(product_info.get('features', [])):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"特性 {i+1}: {feature['title']}"
            tf = slide.placeholders[1].text_frame
            tf.text = feature['description']
            
            if 'points' in feature:
                for point in feature['points']:
                    p = tf.add_paragraph()
                    p.text = f"✓ {point}"
        
        # 总结
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "总结"
        tf = slide.placeholders[1].text_frame
        for benefit in product_info.get('benefits', []):
            p = tf.add_paragraph()
            p.text = f"✓ {benefit}"
        
        prs.save(filepath)
        print(f"✅ 产品演示 PPT 已创建: {filepath}")
        return filepath
    
    @staticmethod
    def create_status_summary(filepath, status_data):
        """创建状态汇总 PPT"""
        prs = Presentation()
        
        # 封面
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = status_data['title']
        slide.placeholders[1].text = f"汇报人: {status_data['reporter']}\n日期: {status_data['date']}"
        
        # 整体状态
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "整体概览"
        
        chart_data = CategoryChartData()
        chart_data.categories = list(status_data['stats'].keys())
        chart_data.add_series('数量', list(status_data['stats'].values()))
        
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(1), Inches(2), Inches(8), Inches(4.5), chart_data
        ).chart
        chart.has_title = False
        
        # 详细信息
        for category, items in status_data.get('details', {}).items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = category
            tf = slide.placeholders[1].text_frame
            for item in items:
                p = tf.add_paragraph()
                p.text = f"• {item}"
        
        prs.save(filepath)
        print(f"✅ 状态汇总 PPT 已创建: {filepath}")
        return filepath


# ==================== 便捷函数 ====================

def quick_create(title):
    """快速创建演示文稿"""
    return PPTProcessor().create().add_title_slide(title)


def create_template_presentation(filepath, slides_data):
    """从数据创建演示文稿"""
    ppt = PPTProcessor()
    ppt.create()
    
    for slide_data in slides_data:
        if slide_data['type'] == 'title':
            ppt.add_title_slide(slide_data['title'], slide_data.get('subtitle', ''))
        elif slide_data['type'] == 'content':
            ppt.add_content_slide(slide_data['title'], slide_data.get('content', ''))
        elif slide_data['type'] == 'bullets':
            ppt.add_content_slide(slide_data['title'])
            ppt.add_bullet_text(slide_data['items'])
        elif slide_data['type'] == 'chart':
            ppt.add_chart(slide_data['title'], slide_data['chart_type'], slide_data['data'])
        elif slide_data['type'] == 'table':
            ppt.add_table(slide_data['data'])
    
    return ppt.save(filepath)


if __name__ == "__main__":
    print("🧪 测试 PPT 处理器")
    
    # 测试创建周报
    print("\n1. 创建测试周报...")
    test_file = "/tmp/test_weekly.pptx"
    PPTProcessor.create_weekly_report(
        test_file,
        {
            'project': '测试项目',
            'date_range': '2024年1月15日 - 1月19日',
            'week': '第3周',
            'completed': ['完成需求分析', '完成UI设计'],
            'in_progress': ['开发进行中', '单元测试'],
            'risks': ['进度可能延期'],
            'next_week': ['继续开发', '准备UAT'],
            'progress': [5, 3, 2]  # 已完成/进行中/待开始
        }
    )
    
    # 测试创建项目报告
    print("\n2. 创建测试项目报告...")
    report_file = "/tmp/test_project.pptx"
    PPTProcessor.create_project_report(
        report_file,
        {
            'title': 'XX系统升级项目',
            'name': 'XX系统',
            'manager': '张三',
            'date': '2024年1月',
            'phase': '开发阶段',
            'progress': 45,
            'team_size': 8,
            'task_stats': {'已完成': 10, '进行中': 5, '待开始': 8},
            'risks': ['技术难点', '人员紧张'],
            'plans': ['完成核心功能', '启动测试']
        }
    )
    
    print("\n3. 测试完成！")
